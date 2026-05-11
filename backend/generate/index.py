import json
import os
import psycopg2
import urllib.request
import urllib.error
import random
import base64
import uuid
import boto3

SCHEMA = "t_p97689468_neural_network_porta"

PHOTO_ROULETTE_PROMPTS = [
    "Утреннее кафе в Париже, золотой свет, дымящийся кофе, пустой стул у окна — фотореалистично",
    "Женщина в красном пальто на фоне осеннего парка, портрет, боке, плёночная фотография",
    "Неоновые вывески Токио ночью, отражения в лужах, одинокая фигура с зонтом",
    "Полярное сияние над заснеженным лесом, звёздное небо, голубой и фиолетовый свет",
    "Заброшенная оранжерея, пробивающийся сквозь стекло свет, буйная растительность",
    "Приморская деревня на рассвете, рыболовные лодки, туман, пастельные тона",
    "Интерьер уютной библиотеки, тёплый свет, стеллажи до потолка, кожаное кресло",
    "Диффузный свет заката над горным хребтом, облака, горизонт в оранжевых тонах",
    "Стрит-фуд рынок в Таиланде, яркие огни, пар от блюд, многолюдность, вечер",
    "Молодая девушка читает книгу под дождём у окна, уличные фонари размыты",
    "Средиземноморские ступени с розовыми кустами, белые стены, синие двери, солнце",
    "Профессиональный фотопортрет в студии, мягкий рассеянный свет, чистый фон",
    "Ночная съёмка города с высоты, огни окон небоскрёбов, длинная выдержка",
    "Завтрак на деревянном столе, акай, свежие ягоды, гранола, минимализм, сверху",
    "Закат на пустынном пляже, следы на песке, волны, золотой час",
    "Цветущая сакура в японском саду, пруд с карпами, деревянный мостик",
    "Уличный художник рисует граффити, краска в воздухе, живые цвета",
    "Стекло с каплями дождя, размытые огни на заднем плане, боке"
]

def get_conn():
    return psycopg2.connect(os.environ["DATABASE_URL"])

def get_user_by_token(cur, token):
    cur.execute(
        f"SELECT id, email, name, is_admin FROM {SCHEMA}.users WHERE password_hash LIKE %s",
        ("%" + token,)
    )
    return cur.fetchone()

def get_s3():
    return boto3.client(
        "s3",
        endpoint_url="https://bucket.poehali.dev",
        aws_access_key_id=os.environ["AWS_ACCESS_KEY_ID"],
        aws_secret_access_key=os.environ["AWS_SECRET_ACCESS_KEY"],
    )

def upload_image_to_s3(image_bytes: bytes, prefix: str = "generated") -> str:
    s3 = get_s3()
    key = f"images/{prefix}/{uuid.uuid4()}.png"
    s3.put_object(Bucket="files", Key=key, Body=image_bytes, ContentType="image/png")
    access_key = os.environ["AWS_ACCESS_KEY_ID"]
    return f"https://cdn.poehali.dev/projects/{access_key}/bucket/{key}"

SIZE_MAP = {
    "square":    (1024, 1024),
    "portrait":  (768, 1024),
    "landscape": (1024, 768),
    "story":     (576, 1024),
    "wide":      (1280, 720),
}

def call_gemini_txt2img(prompt: str, size: str = "square") -> bytes:
    """Генерация изображения через Google Imagen 3 (imagen-3.0-generate-002)"""
    api_key = os.environ.get("IMAGEN_API_KEY", "") or os.environ.get("GEMINI_API_KEY", "")
    if not api_key:
        raise Exception("IMAGEN_API_KEY не настроен")

    ASPECT_MAP = {
        "square":    "1:1",
        "portrait":  "3:4",
        "landscape": "4:3",
        "story":     "9:16",
        "wide":      "16:9",
    }
    aspect = ASPECT_MAP.get(size, "1:1")

    payload = json.dumps({
        "instances": [{"prompt": prompt}],
        "parameters": {
            "sampleCount": 1,
            "aspectRatio": aspect,
            "safetyFilterLevel": "block_few",
            "personGeneration": "allow_adult",
        }
    }).encode()

    url = f"https://generativelanguage.googleapis.com/v1/models/imagen-3.0-generate-002:predict?key={api_key}"
    req = urllib.request.Request(url, data=payload, method="POST")
    req.add_header("Content-Type", "application/json")

    try:
        with urllib.request.urlopen(req, timeout=120) as resp:
            data = json.loads(resp.read())
    except urllib.error.HTTPError as e:
        err = e.read().decode("utf-8", errors="ignore")
        print(f"[imagen3] HTTP {e.code}: {err[:500]}")
        raise Exception(f"Imagen 3 HTTP {e.code}: {err[:300]}")

    predictions = data.get("predictions", [])
    if not predictions:
        raise Exception(f"Imagen 3: пустой ответ — {str(data)[:200]}")
    b64 = predictions[0].get("bytesBase64Encoded", "")
    if not b64:
        raise Exception(f"Imagen 3: нет изображения в ответе")
    return base64.b64decode(b64)


def call_pollinations_txt2img(prompt: str, size: str = "square") -> bytes:
    """Резервный генератор через Pollinations (используется если Gemini недоступен)"""
    import urllib.parse
    w, h = SIZE_MAP.get(size, (1024, 1024))
    encoded = urllib.parse.quote(prompt)
    url = f"https://image.pollinations.ai/prompt/{encoded}?width={w}&height={h}&model=flux&nologo=true&enhance=false"
    req = urllib.request.Request(url, method="GET")
    req.add_header("User-Agent", "Mozilla/5.0")
    with urllib.request.urlopen(req, timeout=120) as resp:
        return resp.read()


def call_gpt_image(prompt: str, size: str = "square") -> bytes:
    """Генерация изображения через GPT Image 1 (gpt-image-1)"""
    api_key = os.environ.get("OPENAI_API_KEY", "")
    if not api_key:
        raise Exception("OPENAI_API_KEY не настроен")

    SIZE_OPENAI = {
        "square":    "1024x1024",
        "portrait":  "1024x1536",
        "landscape": "1536x1024",
        "story":     "1024x1536",
        "wide":      "1536x1024",
    }
    img_size = SIZE_OPENAI.get(size, "1024x1024")

    payload = json.dumps({
        "model": "gpt-image-1",
        "prompt": prompt,
        "n": 1,
        "size": img_size,
        "quality": "standard",
        "output_format": "png",
    }).encode()

    req = urllib.request.Request("https://api.openai.com/v1/images/generations", data=payload, method="POST")
    req.add_header("Content-Type", "application/json")
    req.add_header("Authorization", f"Bearer {api_key}")

    try:
        with urllib.request.urlopen(req, timeout=120) as resp:
            data = json.loads(resp.read())
    except urllib.error.HTTPError as e:
        err = e.read().decode("utf-8", errors="ignore")
        print(f"[gpt-image] HTTP {e.code}: {err[:500]}")
        raise Exception(f"GPT Image HTTP {e.code}: {err[:300]}")

    b64 = data.get("data", [{}])[0].get("b64_json", "")
    if not b64:
        raise Exception(f"GPT Image: нет изображения в ответе — {str(data)[:200]}")
    return base64.b64decode(b64)


def call_huggingface_txt2img(prompt: str, size: str = "square") -> bytes:
    """Генерация изображения через Hugging Face FLUX.1-dev"""
    token = os.environ.get("HUGGINGFACE_TOKEN", "")
    if not token:
        raise Exception("HUGGINGFACE_TOKEN не настроен")

    HF_SIZE_MAP = {
        "square":    (1024, 1024),
        "portrait":  (768, 1024),
        "landscape": (1024, 768),
        "story":     (576, 1024),
        "wide":      (1280, 720),
    }
    w, h = HF_SIZE_MAP.get(size, (1024, 1024))

    payload = json.dumps({
        "inputs": prompt,
        "parameters": {"num_inference_steps": 30, "guidance_scale": 7.5},
        "options": {"wait_for_model": True, "use_cache": False}
    }).encode()

    url = "https://api-inference.huggingface.co/models/stabilityai/stable-diffusion-xl-base-1.0"
    req = urllib.request.Request(url, data=payload, method="POST")
    req.add_header("Authorization", f"Bearer {token}")
    req.add_header("Content-Type", "application/json")

    try:
        with urllib.request.urlopen(req, timeout=120) as resp:
            data = resp.read()
            content_type = resp.headers.get("Content-Type", "")
            if content_type.startswith("image/"):
                return data
            result = json.loads(data)
            if isinstance(result, dict) and result.get("image"):
                return base64.b64decode(result["image"])
            raise Exception(f"HuggingFace неожиданный ответ: {str(result)[:200]}")
    except urllib.error.HTTPError as e:
        err = e.read().decode("utf-8", errors="ignore")
        print(f"[huggingface] HTTP {e.code}: {err[:500]}")
        raise Exception(f"HuggingFace HTTP {e.code}: {err[:300]}")


def call_stability_txt2img(prompt: str, size: str = "square") -> bytes:
    """Генерация изображения через Stability AI (stable-diffusion-3-5-large)"""
    api_key = os.environ.get("STABILITY_API_KEY", "")
    if not api_key:
        raise Exception("STABILITY_API_KEY не настроен")

    ASPECT_MAP = {
        "square":    "1:1",
        "portrait":  "3:4",
        "landscape": "4:3",
        "story":     "9:16",
        "wide":      "16:9",
    }
    aspect = ASPECT_MAP.get(size, "1:1")

    import uuid as _uuid
    boundary = _uuid.uuid4().hex
    def field(name, value):
        return (f"--{boundary}\r\nContent-Disposition: form-data; name=\"{name}\"\r\n\r\n{value}\r\n").encode()
    body = field("prompt", prompt) + field("aspect_ratio", aspect) + field("output_format", "png") + field("model", "sd3.5-large") + f"--{boundary}--\r\n".encode()

    url = "https://api.stability.ai/v2beta/stable-image/generate/sd3"
    req = urllib.request.Request(url, data=body, method="POST")
    req.add_header("Authorization", f"Bearer {api_key}")
    req.add_header("Accept", "image/*")
    req.add_header("Content-Type", f"multipart/form-data; boundary={boundary}")

    try:
        with urllib.request.urlopen(req, timeout=120) as resp:
            return resp.read()
    except urllib.error.HTTPError as e:
        err = e.read().decode("utf-8", errors="ignore")
        print(f"[stability] HTTP {e.code}: {err[:500]}")
        raise Exception(f"Stability AI HTTP {e.code}: {err[:300]}")


def generate_image_with_fallback(prompt: str, size: str = "square") -> bytes:
    """Цепочка: Stability AI SD3.5 → Imagen 3"""
    try:
        print(f"[image] Stability AI SD3.5: {prompt[:80]}")
        return call_stability_txt2img(prompt, size)
    except Exception as e:
        print(f"[image] Stability failed ({e}), fallback → Imagen 3")
    return call_gemini_txt2img(prompt, size)

def resize_image(image_bytes: bytes, max_size: int = 512) -> bytes:
    """Сжимает изображение до max_size по большей стороне через PIL"""
    from PIL import Image
    import io
    img = Image.open(io.BytesIO(image_bytes)).convert("RGB")
    w, h = img.size
    if max(w, h) > max_size:
        ratio = max_size / max(w, h)
        img = img.resize((int(w * ratio), int(h * ratio)), Image.LANCZOS)
    buf = io.BytesIO()
    img.save(buf, format="JPEG", quality=85)
    return buf.getvalue()

def call_cloudflare_img2img(image_bytes: bytes, prompt: str) -> bytes:
    """Редактирование фото через Cloudflare Workers AI (flux-1-schnell-img2img)"""
    import io
    from PIL import Image

    account_id = os.environ.get("CLOUDFLARE_ACCOUNT_ID", "")
    api_token = os.environ.get("CLOUDFLARE_API_TOKEN", "")
    if not account_id or not api_token:
        raise Exception("CLOUDFLARE_ACCOUNT_ID или CLOUDFLARE_API_TOKEN не настроены")

    # Сжимаем до 512x512 (лимит CF img2img)
    img = Image.open(io.BytesIO(image_bytes)).convert("RGB")
    img.thumbnail((512, 512), Image.LANCZOS)
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    image_array = list(buf.getvalue())

    payload = json.dumps({
        "prompt": prompt,
        "image": image_array,
        "strength": 0.75,
        "num_steps": 20,
        "guidance": 7.5,
    }).encode()

    url = f"https://api.cloudflare.com/client/v4/accounts/{account_id}/ai/run/@cf/runwayml/stable-diffusion-v1-5-img2img"
    req = urllib.request.Request(url, data=payload, method="POST")
    req.add_header("Authorization", f"Bearer {api_token}")
    req.add_header("Content-Type", "application/json")

    try:
        with urllib.request.urlopen(req, timeout=120) as resp:
            raw = resp.read()
    except urllib.error.HTTPError as e:
        err_body = e.read().decode("utf-8", errors="ignore")
        print(f"[cloudflare img2img] HTTP {e.code}: {err_body[:1000]}")
        raise Exception(f"Cloudflare HTTP {e.code}: {err_body[:300]}")

    # stable-diffusion-v1-5-img2img возвращает сырые байты PNG
    if raw[:4] == b'\x89PNG' or raw[:2] == b'\xff\xd8':
        print(f"[cloudflare img2img] got raw image: {len(raw)} bytes")
        return raw

    # Или JSON с base64
    try:
        result = json.loads(raw)
        print(f"[cloudflare img2img] response: {str(result)[:300]}")
        if result.get("result") and isinstance(result["result"], str):
            return base64.b64decode(result["result"])
        if result.get("result", {}).get("image"):
            return base64.b64decode(result["result"]["image"])
        raise Exception(f"Cloudflare неожиданный ответ: {str(result)[:200]}")
    except (json.JSONDecodeError, KeyError):
        raise Exception(f"Cloudflare вернул нераспознанный формат: {len(raw)} bytes")

def gemini_generate_slides(topic: str, slides_count: int, ai_text: bool) -> list:
    """Генерирует структуру слайдов карусели через OpenRouter (gpt-4o-mini)"""
    api_key = os.environ.get("OPENROUTER_API_KEY", "")
    if not api_key:
        raise Exception("OPENROUTER_API_KEY не настроен")

    text_instruction = (
        "Для каждого слайда придумай яркий заголовок И полный готовый текст для слайда (2-4 строки, конкретно и ёмко)."
        if ai_text else
        "Для каждого слайда придумай яркий заголовок и краткое описание содержимого (подсказку автору)."
    )
    prompt = (
        f"Создай структуру карусели из {slides_count} слайдов для Instagram на тему: '{topic}'.\n"
        f"{text_instruction}\n"
        f"Для каждого слайда придумай описание визуала на английском (5-10 слов, что нарисовать на фоне).\n"
        f"Верни ТОЛЬКО JSON-массив без markdown:\n"
        f'[{{"slide": 1, "title": "...", "text": "...", "visual_prompt": "..."}}]\n'
        f"Первый слайд — цепляющий крючок. Последний — призыв к действию."
    )
    payload = json.dumps({
        "model": "openai/gpt-4o-mini",
        "messages": [
            {"role": "system", "content": "Ты — создатель вирусного контента для Instagram. Отвечай ТОЛЬКО валидным JSON-массивом без пояснений и markdown."},
            {"role": "user", "content": prompt}
        ],
        "max_tokens": 4000
    }).encode()
    req = urllib.request.Request("https://openrouter.ai/api/v1/chat/completions", data=payload, method="POST")
    req.add_header("Content-Type", "application/json")
    req.add_header("Authorization", f"Bearer {api_key}")
    req.add_header("HTTP-Referer", "https://neuralai.poehali.dev")
    try:
        with urllib.request.urlopen(req, timeout=60) as resp:
            data = json.loads(resp.read())
    except urllib.error.HTTPError as e:
        err = e.read().decode("utf-8", errors="ignore")
        raise Exception(f"OpenRouter error {e.code}: {err[:300]}")
    raw = data["choices"][0]["message"]["content"]
    clean = raw.strip()
    if clean.startswith("```"):
        parts = clean.split("```")
        clean = parts[1] if len(parts) > 1 else clean
        if clean.startswith("json"):
            clean = clean[4:]
        clean = clean.strip()
    return json.loads(clean)


def generate_slide_image_pollinations(visual_prompt: str, style_prompt: str) -> bytes:
    """Генерирует изображение слайда через Gemini Imagen 3 (с fallback на Pollinations)"""
    full_prompt = f"{visual_prompt}, {style_prompt}, high quality, instagram post, no text overlay"
    return generate_image_with_fallback(full_prompt, size="square")


def generate_slide_image_with_photo(visual_prompt: str, style_prompt: str, user_image_b64: str) -> bytes:
    """Генерирует изображение слайда с пользовательским фото через CF img2img"""
    from PIL import Image
    import io
    account_id = os.environ.get("CLOUDFLARE_ACCOUNT_ID", "")
    api_token = os.environ.get("CLOUDFLARE_API_TOKEN", "")
    if not account_id or not api_token:
        raise Exception("Cloudflare не настроен")
    img_bytes = base64.b64decode(user_image_b64)
    img = Image.open(io.BytesIO(img_bytes)).convert("RGB")
    img.thumbnail((512, 512), Image.LANCZOS)
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    image_array = list(buf.getvalue())
    full_prompt = f"{visual_prompt}, {style_prompt}, instagram carousel, no text"
    payload = json.dumps({
        "prompt": full_prompt,
        "image": image_array,
        "strength": 0.6,
        "num_steps": 20,
        "guidance": 7.5,
    }).encode()
    url = f"https://api.cloudflare.com/client/v4/accounts/{account_id}/ai/run/@cf/runwayml/stable-diffusion-v1-5-img2img"
    req = urllib.request.Request(url, data=payload, method="POST")
    req.add_header("Authorization", f"Bearer {api_token}")
    req.add_header("Content-Type", "application/json")
    try:
        with urllib.request.urlopen(req, timeout=120) as resp:
            raw = resp.read()
    except urllib.error.HTTPError as e:
        err = e.read().decode("utf-8", errors="ignore")
        raise Exception(f"CF error {e.code}: {err[:200]}")
    if raw[:4] == b'\x89PNG' or raw[:2] == b'\xff\xd8':
        return raw
    result = json.loads(raw)
    if result.get("result", {}).get("image"):
        return base64.b64decode(result["result"]["image"])
    raise Exception(f"CF img2img неожиданный ответ")


def generate_text_with_openrouter(prompt: str, system: str = "") -> str:
    api_key = os.environ.get("OPENROUTER_API_KEY", "")
    if not api_key:
        return "API ключ не настроен. Обратитесь к администратору."

    messages = []
    if system:
        messages.append({"role": "system", "content": system})
    messages.append({"role": "user", "content": prompt})

    payload = json.dumps({
        "model": "openai/gpt-4o-mini",
        "messages": messages,
        "max_tokens": 2000
    }).encode()

    req = urllib.request.Request(
        "https://openrouter.ai/api/v1/chat/completions",
        data=payload, method="POST"
    )
    req.add_header("Content-Type", "application/json")
    req.add_header("Authorization", f"Bearer {api_key}")
    req.add_header("HTTP-Referer", "https://neuralai.poehali.dev")

    with urllib.request.urlopen(req, timeout=60) as resp:
        data = json.loads(resp.read())
        return data["choices"][0]["message"]["content"]

def handler(event: dict, context) -> dict:
    """Генерация контента: изображения по промту, редактирование фото, тексты, карусели, сценарии"""
    headers = {
        "Access-Control-Allow-Origin": "*",
        "Access-Control-Allow-Methods": "GET, POST, PUT, DELETE, OPTIONS",
        "Access-Control-Allow-Headers": "Content-Type, X-User-Id, X-Auth-Token, X-Session-Id, X-Authorization, Authorization",
        "Content-Type": "application/json"
    }

    if event.get("httpMethod") == "OPTIONS":
        return {"statusCode": 200, "headers": headers, "body": ""}

    method = event.get("httpMethod", "GET")
    path = event.get("path", "/")
    body = {}
    if event.get("body"):
        try:
            body = json.loads(event["body"])
        except Exception:
            pass

    if method == "GET":
        prompt = random.choice(PHOTO_ROULETTE_PROMPTS)
        return {"statusCode": 200, "headers": headers, "body": json.dumps({"prompt": prompt})}

    if method != "POST":
        return {"statusCode": 405, "headers": headers, "body": json.dumps({"error": "Метод не разрешён"})}

    action = body.get("action", "")

    if action == "image-gen":
        prompt = body.get("prompt", "")
        style = body.get("style", "")
        size = body.get("size", "square")
        if not prompt:
            return {"statusCode": 400, "headers": headers, "body": json.dumps({"error": "Укажите описание изображения"})}

        full_prompt = f"{prompt}, {style}" if style else prompt
        full_prompt += ", high quality, 8k, photorealistic"

        try:
            image_bytes = generate_image_with_fallback(full_prompt, size)
            cdn_url = upload_image_to_s3(image_bytes, prefix="generated")
            return {"statusCode": 200, "headers": headers, "body": json.dumps({"image_url": cdn_url, "prompt": full_prompt})}
        except Exception as e:
            error_msg = str(e)
            print(f"[image-gen error] {error_msg}")
            return {"statusCode": 500, "headers": headers, "body": json.dumps({"error": f"Ошибка генерации: {error_msg}"})}

    elif action == "image-edit":
        prompt = body.get("prompt", "")
        image_b64 = body.get("image_base64", "")
        size = body.get("size", "square")
        if not prompt or not image_b64:
            return {"statusCode": 400, "headers": headers, "body": json.dumps({"error": "Укажите изображение и описание изменений"})}

        try:
            image_bytes = base64.b64decode(image_b64)
        except Exception:
            return {"statusCode": 400, "headers": headers, "body": json.dumps({"error": "Неверный формат изображения"})}

        try:
            result_bytes = call_cloudflare_img2img(image_bytes, prompt)
            cdn_url = upload_image_to_s3(result_bytes, prefix="edited")
            return {"statusCode": 200, "headers": headers, "body": json.dumps({"image_url": cdn_url, "prompt": prompt})}
        except Exception as e:
            error_msg = str(e)
            print(f"[image-edit error] {error_msg}")
            return {"statusCode": 500, "headers": headers, "body": json.dumps({"error": f"Ошибка редактирования: {error_msg}"})}

    elif action == "post":
        topic = body.get("topic", "")
        platform = body.get("platform", "Instagram")
        tone = body.get("tone", "вовлекающий")
        length = body.get("length", "Средний")
        language = body.get("language", "Русский")
        emoji_style = body.get("emoji_style", "Умеренно")

        length_map = {"Короткий": "до 500 символов", "Средний": "500-1000 символов", "Длинный": "1000-2000 символов"}
        emoji_map = {"Много эмодзи 🔥": "используй много эмодзи в каждом предложении", "Умеренно": "используй эмодзи умеренно", "Без эмодзи": "без эмодзи вообще"}
        lang_instruction = "Пиши на английском языке." if language == "English" else "Пиши на русском языке."

        system = f"Ты — профессиональный копирайтер для социальных сетей. {lang_instruction}"
        prompt = (
            f"Напиши {tone} пост для {platform} на тему: '{topic}'.\n"
            f"Длина: {length_map.get(length, '500-1000 символов')}.\n"
            f"Эмодзи: {emoji_map.get(emoji_style, 'умеренно')}.\n"
            f"Добавь хэштеги и призыв к действию."
        )
        result = generate_text_with_openrouter(prompt, system)
        return {"statusCode": 200, "headers": headers, "body": json.dumps({"result": result, "type": "post"})}

    elif action == "hashtags":
        topic = body.get("topic", "")
        platform = body.get("platform", "Instagram")
        if not topic:
            return {"statusCode": 400, "headers": headers, "body": json.dumps({"error": "Укажите тему"})}
        system = "Ты — эксперт по SMM и хэштегам. Отвечай ТОЛЬКО валидным JSON без markdown."
        prompt = (
            f"Подбери 30 хэштегов для {platform} по теме: '{topic}'.\n"
            f"Раздели на три группы по частотности:\n"
            f"- high: 10 высокочастотных (миллионы публикаций)\n"
            f"- medium: 10 среднечастотных (100к-1млн публикаций)\n"
            f"- low: 10 низкочастотных (до 100к, нишевые)\n"
            f"Верни JSON: {{\"high\": [...], \"medium\": [...], \"low\": [...]}}\n"
            f"Хэштеги с # в начале, на русском или английском (как принято в нише)."
        )
        raw = generate_text_with_openrouter(prompt, system)
        try:
            clean = raw.strip()
            if clean.startswith("```"):
                parts = clean.split("```")
                clean = parts[1] if len(parts) > 1 else clean
                if clean.startswith("json"): clean = clean[4:]
                clean = clean.strip()
            data_h = json.loads(clean)
            all_tags = data_h.get("high", []) + data_h.get("medium", []) + data_h.get("low", [])
            data_h["all"] = all_tags
            return {"statusCode": 200, "headers": headers, "body": json.dumps(data_h)}
        except Exception:
            return {"statusCode": 200, "headers": headers, "body": json.dumps({"high": [], "medium": [], "low": [], "all": [], "raw": raw})}

    elif action == "bio-generator":
        niche = body.get("niche", "")
        keywords = body.get("keywords", "")
        contact = body.get("contact", "")
        platform = body.get("platform", "Instagram")
        tone = body.get("tone", "экспертный")
        if not niche:
            return {"statusCode": 400, "headers": headers, "body": json.dumps({"error": "Укажите нишу"})}
        contact_str = f"Контакт/ссылка: {contact}." if contact else ""
        system = "Ты — эксперт по личному брендингу в соцсетях. Отвечай ТОЛЬКО валидным JSON без markdown."
        prompt = (
            f"Создай 3 варианта шапки профиля для {platform}.\n"
            f"Ниша: {niche}. Ключевые слова: {keywords}. {contact_str} Тон: {tone}.\n"
            f"Каждый вариант: 2-4 строки, эмодзи, конкретная польза для подписчика.\n"
            f"Учти лимит символов {platform}: Instagram 150 симв, TikTok 80 симв, Telegram 70 симв, VK 200 симв.\n"
            f"Верни JSON: {{\"variants\": [\"вариант 1\", \"вариант 2\", \"вариант 3\"]}}"
        )
        raw = generate_text_with_openrouter(prompt, system)
        try:
            clean = raw.strip()
            if clean.startswith("```"):
                parts = clean.split("```")
                clean = parts[1] if len(parts) > 1 else clean
                if clean.startswith("json"): clean = clean[4:]
                clean = clean.strip()
            data_b = json.loads(clean)
            return {"statusCode": 200, "headers": headers, "body": json.dumps(data_b)}
        except Exception:
            return {"statusCode": 200, "headers": headers, "body": json.dumps({"variants": [raw]})}

    elif action == "repurpose":
        text = body.get("text", "")
        formats = body.get("formats", [])
        if not text or not formats:
            return {"statusCode": 400, "headers": headers, "body": json.dumps({"error": "Укажите текст и форматы"})}
        results = {}
        format_prompts = {
            "Пост Instagram": f"Переупакуй этот текст в продающий пост для Instagram (до 1000 символов, эмодзи, хэштеги, призыв к действию):\n\n{text}",
            "Карусель": f"Переупакуй этот текст в структуру карусели для Instagram (5-7 слайдов, для каждого: заголовок + 2-3 тезиса):\n\n{text}",
            "Сторис": f"Переупакуй этот текст в 3 части для сторис (каждая часть — 1-2 предложения, клиффхэнгер между частями):\n\n{text}",
            "Email-рассылка": f"Переупакуй этот текст в email-рассылку (тема письма, приветствие, основной текст, призыв к действию, подпись):\n\n{text}",
            "Заголовок для Дзена": f"Придумай 5 цепляющих заголовков для Яндекс Дзена на основе этого текста (с числами, вопросами, интригой):\n\n{text}",
            "Thread для Telegram": f"Переупакуй этот текст в тред для Telegram (5-7 коротких сообщений пронумерованных 1/7, 2/7 и т.д.):\n\n{text}",
        }
        system = "Ты — профессиональный контент-маркетолог. Пиши на русском языке, адаптируй контент под формат."
        for fmt in formats:
            if fmt in format_prompts:
                results[fmt] = generate_text_with_openrouter(format_prompts[fmt], system)
        return {"statusCode": 200, "headers": headers, "body": json.dumps({"results": results})}

    elif action == "avatar-gen":
        image_b64 = body.get("image_b64", "")
        style_prompt = body.get("style_prompt", "professional portrait")
        if not image_b64:
            return {"statusCode": 400, "headers": headers, "body": json.dumps({"error": "Загрузите фото"})}
        try:
            img_bytes = base64.b64decode(image_b64)
            result_bytes = generate_slide_image_with_photo(style_prompt, "high quality portrait, detailed", image_b64)
            cdn_url = upload_image_to_s3(result_bytes, prefix="avatars")
            return {"statusCode": 200, "headers": headers, "body": json.dumps({"image_url": cdn_url})}
        except Exception as e:
            print(f"[avatar-gen error] {e}")
            return {"statusCode": 500, "headers": headers, "body": json.dumps({"error": f"Ошибка генерации: {str(e)}"})}

    elif action == "stories-gen":
        topic = body.get("topic", "")
        main_text = body.get("main_text", "")
        sub_text = body.get("sub_text", "")
        bg_prompt = body.get("bg_prompt", "vibrant gradient background")
        if not main_text:
            return {"statusCode": 400, "headers": headers, "body": json.dumps({"error": "Введите текст для Stories"})}
        try:
            # Генерируем фон через Gemini Imagen 3 (с fallback на Pollinations)
            full_prompt = f"{bg_prompt}, vertical format, instagram story background, no text, no letters, abstract"
            img_bytes = generate_image_with_fallback(full_prompt, size="story")
            cdn_url = upload_image_to_s3(img_bytes, prefix="stories")
            return {"statusCode": 200, "headers": headers, "body": json.dumps({"image_url": cdn_url, "main_text": main_text, "sub_text": sub_text})}
        except Exception as e:
            print(f"[stories-gen error] {e}")
            return {"statusCode": 500, "headers": headers, "body": json.dumps({"error": f"Ошибка: {str(e)}"})}

    elif action == "brand-kit-analysis":
        brand_name = body.get("brandName", "")
        niche = body.get("niche", "")
        audience = body.get("targetAudience", "")
        tone = body.get("tone", "")
        keywords = body.get("keywords", "")
        usp = body.get("usp", "")
        colors = body.get("mainColors", "")
        system = "Ты — эксперт по личному брендингу и маркетингу в соцсетях. Пиши структурированно и конкретно."
        prompt = (
            f"Проведи анализ бренда и дай конкретные рекомендации:\n"
            f"Бренд: {brand_name}\nНиша: {niche}\nАудитория: {audience}\n"
            f"Тон: {tone}\nКлючевые слова: {keywords}\nУТП: {usp}\nЦвета: {colors}\n\n"
            f"Дай анализ по разделам:\n"
            f"1. Сильные стороны бренда\n"
            f"2. Что улучшить в позиционировании\n"
            f"3. Рекомендуемые форматы контента\n"
            f"4. Ключевые темы для постов (10 идей)\n"
            f"5. Советы по tone of voice\n"
            f"6. Потенциальные точки роста"
        )
        result = generate_text_with_openrouter(prompt, system)
        return {"statusCode": 200, "headers": headers, "body": json.dumps({"result": result})}

    elif action == "ai-assistant":
        message = body.get("message", "")
        history = body.get("history", [])
        if not message:
            return {"statusCode": 400, "headers": headers, "body": json.dumps({"error": "Сообщение пустое"})}

        TOOL_MAP = {
            "пост": ("/tools/post", "Написать пост", "Wand2"),
            "карусель": ("/tools/carousel", "Создать карусель", "LayoutTemplate"),
            "хэштег": ("/tools/hashtags", "Подобрать хэштеги", "Hash"),
            "контент-план": ("/tools/content-plan", "Контент-план", "CalendarDays"),
            "изображен": ("/tools/image-gen", "Сгенерировать фото", "ImagePlus"),
            "презентац": ("/tools/presentation", "Создать презентацию", "Presentation"),
            "сценарий": ("/tools/scenario", "Написать сценарий", "Clapperboard"),
            "скрипт": ("/tools/sale-script", "Скрипт продаж", "MessageSquare"),
            "email": ("/tools/email", "Email-копирайтер", "Mail"),
            "шапк": ("/tools/bio", "Шапка профиля", "UserCircle"),
            "кейс": ("/tools/case", "Генератор кейсов", "Trophy"),
            "конкурент": ("/tools/competitor", "Анализ конкурентов", "Search"),
            "stories": ("/tools/stories", "Генератор Stories", "Smartphone"),
            "аватар": ("/tools/avatar", "ИИ-аватар", "UserCircle2"),
            "гайд": ("/tools/guide", "Создать гайд", "BookOpen"),
            "воронк": ("/tools/funnel", "Воронка продаж", "TrendingUp"),
        }

        system = (
            "Ты — дружелюбный ИИ-ассистент платформы Neural AI для создания контента. "
            "Отвечай кратко (2-4 предложения), по-русски, с конкретной помощью. "
            "Если пользователь просит создать контент — выполни задачу прямо здесь (напиши пост, хэштеги, идеи и т.д.), "
            "а не просто отправь куда-то. Будь полезным и тёплым."
        )
        messages_list = []
        for h in history[-6:]:
            messages_list.append({"role": h.get("role", "user"), "content": h.get("text", "")})
        messages_list.append({"role": "user", "content": message})

        payload = json.dumps({
            "model": "openai/gpt-4o-mini",
            "messages": [{"role": "system", "content": system}] + messages_list,
            "max_tokens": 600
        }).encode()
        req = urllib.request.Request("https://openrouter.ai/api/v1/chat/completions", data=payload, method="POST")
        req.add_header("Content-Type", "application/json")
        req.add_header("Authorization", f"Bearer {os.environ.get('OPENROUTER_API_KEY', '')}")
        req.add_header("HTTP-Referer", "https://neuralai.poehali.dev")

        try:
            with urllib.request.urlopen(req, timeout=30) as resp:
                data_r = json.loads(resp.read())
            reply = data_r["choices"][0]["message"]["content"]
        except Exception as e:
            reply = "Извини, произошла ошибка. Попробуй ещё раз!"

        # Определяем подходящие инструменты по ключевым словам
        msg_lower = message.lower()
        actions = []
        for keyword, (href, label, icon) in TOOL_MAP.items():
            if keyword in msg_lower and len(actions) < 3:
                actions.append({"label": label, "href": href, "icon": icon})

        return {"statusCode": 200, "headers": headers, "body": json.dumps({
            "reply": reply,
            "actions": actions if actions else None
        })}

    elif action == "naming":
        business_type = body.get("business_type", "")
        niche = body.get("niche", "")
        style = body.get("style", "Дружелюбный")
        target_audience = body.get("target_audience", "")
        with_slogan = body.get("with_slogan", True)
        if not business_type:
            return {"statusCode": 400, "headers": headers, "body": json.dumps({"error": "Укажите тип бизнеса"})}
        slogan_part = "и слоган (1 предложение)" if with_slogan else ""
        system = "Ты — брендинг-эксперт и копирайтер. Отвечай ТОЛЬКО валидным JSON-массивом без markdown."
        prompt = (
            f"Придумай 6 вариантов названия {slogan_part} и краткое описание для:\n"
            f"Бизнес: {business_type}\nНиша: {niche}\nАудитория: {target_audience}\nСтиль: {style}\n\n"
            f"Верни JSON-массив: [{{\"name\": \"...\", \"slogan\": \"...\", \"description\": \"почему это название работает\"}}]\n"
            f"Названия должны быть оригинальными, запоминающимися, подходящими для домена и соцсетей."
        )
        result = generate_text_with_openrouter(prompt, system)
        return {"statusCode": 200, "headers": headers, "body": json.dumps({"result": result, "type": "naming"})}

    elif action == "ad-copy":
        product = body.get("product", "")
        audience = body.get("audience", "")
        platform = body.get("platform", "Яндекс.Директ")
        goal = body.get("goal", "Продажи")
        budget = body.get("budget", "")
        if not product:
            return {"statusCode": 400, "headers": headers, "body": json.dumps({"error": "Укажите продукт"})}
        budget_str = f"Бюджет: {budget}." if budget else ""
        system = "Ты — эксперт по контекстной и таргетированной рекламе. Отвечай ТОЛЬКО валидным JSON-массивом без markdown."
        prompt = (
            f"Создай 4 рекламных объявления для {platform}.\n"
            f"Продукт: {product}\nАудитория: {audience}\nЦель: {goal}\n{budget_str}\n\n"
            f"Учти ограничения {platform}: для Яндекс.Директ заголовок до 56 символов, текст до 81 символа.\n"
            f"Верни JSON-массив: [{{\"title\": \"...\", \"text\": \"...\", \"cta\": \"Кнопка призыва\", \"tip\": \"Совет по этому объявлению\"}}]\n"
            f"Каждое объявление должно быть цепляющим и релевантным аудитории."
        )
        result = generate_text_with_openrouter(prompt, system)
        return {"statusCode": 200, "headers": headers, "body": json.dumps({"result": result, "type": "ad-copy"})}

    elif action == "comments":
        post_topic = body.get("post_topic", "")
        account_niche = body.get("account_niche", "")
        comment_style = body.get("comment_style", "Вопрос")
        count = body.get("count", 10)
        if not post_topic:
            return {"statusCode": 400, "headers": headers, "body": json.dumps({"error": "Укажите тему поста"})}
        system = "Ты — эксперт по продвижению в соцсетях. Отвечай ТОЛЬКО валидным JSON-массивом без markdown."
        prompt = (
            f"Придумай {count} естественных комментариев для прогрева под пост на тему: '{post_topic}'.\n"
            f"Ниша аккаунта: {account_niche}\nСтиль: {comment_style}\n\n"
            f"Комментарии должны быть разными по длине (от 5 до 30 слов), живыми, без рекламы.\n"
            f"Верни JSON-массив строк: [\"комментарий 1\", \"комментарий 2\", ...]"
        )
        result = generate_text_with_openrouter(prompt, system)
        return {"statusCode": 200, "headers": headers, "body": json.dumps({"result": result, "type": "comments"})}

    elif action == "scenario-convert":
        source_text = body.get("source_text", "")
        platform = body.get("platform", "Reels")
        duration = body.get("duration", "60 секунд")
        if not source_text:
            return {"statusCode": 400, "headers": headers, "body": json.dumps({"error": "Вставьте текст для конвертации"})}
        system = "Ты — профессиональный сценарист видео-контента. Пиши на русском языке."
        prompt = (
            f"Конвертируй этот текст в видео-сценарий для {platform}, длительность {duration}:\n\n{source_text}\n\n"
            f"Адаптируй под видео-формат: добавь крючок в начало, разбей на сцены с таймкодами, "
            f"сделай динамичный монтаж, добавь призыв к действию в конце. "
            f"Укажи: [0:00] Крючок, [0:05] Основная часть (по сценам), [конец] CTA, советы по съёмке."
        )
        result = generate_text_with_openrouter(prompt, system)
        return {"statusCode": 200, "headers": headers, "body": json.dumps({"result": result, "type": "scenario"})}

    elif action == "sale-script":
        product = body.get("product", "")
        audience = body.get("audience", "")
        objections = body.get("objections", "")
        channel = body.get("channel", "Переписка")
        script_type = body.get("script_type", "Холодный контакт")
        if not product:
            return {"statusCode": 400, "headers": headers, "body": json.dumps({"error": "Укажите продукт"})}
        system = "Ты — эксперт по продажам с 10-летним опытом. Пиши практичные скрипты продаж на русском языке."
        prompt = (
            f"Напиши детальный скрипт продаж для канала '{channel}', тип: '{script_type}'.\n"
            f"Продукт/услуга: {product}\n"
            f"Целевая аудитория: {audience}\n"
            f"Основные возражения: {objections}\n\n"
            f"Структура скрипта:\n"
            f"1. Открытие/приветствие (с цеплялкой)\n"
            f"2. Выявление потребности (3-5 вопросов)\n"
            f"3. Презентация решения (на языке выгод)\n"
            f"4. Отработка возражений (конкретные фразы для каждого)\n"
            f"5. Закрытие сделки (2-3 варианта)\n"
            f"6. Скрипт для случая отказа\n"
            f"Дай конкретные фразы и реплики, не общие рекомендации."
        )
        result = generate_text_with_openrouter(prompt, system)
        return {"statusCode": 200, "headers": headers, "body": json.dumps({"result": result, "type": "sale-script"})}

    elif action == "email-copy":
        purpose = body.get("purpose", "Продающее")
        product = body.get("product", "")
        audience = body.get("audience", "")
        tone = body.get("tone", "Дружелюбный")
        chain = body.get("chain", False)
        if not product:
            return {"statusCode": 400, "headers": headers, "body": json.dumps({"error": "Укажите продукт или тему"})}
        system = "Ты — профессиональный email-маркетолог. Пиши письма которые открывают и читают. Русский язык."
        if chain:
            prompt = (
                f"Напиши цепочку из 3 писем для email-рассылки.\n"
                f"Тип: {purpose}. Продукт/тема: {product}. Аудитория: {audience}. Тон: {tone}.\n"
                f"Письмо 1 (день 0): знакомство, ценность без продажи\n"
                f"Письмо 2 (день 2): кейс/история/польза, мягкое предложение\n"
                f"Письмо 3 (день 4): ограниченное предложение, призыв к действию\n"
                f"Для каждого письма: тема письма, превью (preheader), полный текст, CTA-кнопка."
            )
        else:
            prompt = (
                f"Напиши {purpose.lower()} письмо для email-рассылки.\n"
                f"Продукт/тема: {product}. Аудитория: {audience}. Тон: {tone}.\n"
                f"Структура: тема письма (с emoji), preheader, приветствие, основной текст, "
                f"призыв к действию (CTA), подпись. Письмо должно быть читабельным и цепляющим."
            )
        result = generate_text_with_openrouter(prompt, system)
        return {"statusCode": 200, "headers": headers, "body": json.dumps({"result": result, "type": "email-copy"})}

    elif action == "competitor-analysis":
        niche = body.get("niche", "")
        competitors = body.get("competitors", "")
        my_strengths = body.get("my_strengths", "")
        analysis_type = body.get("analysis_type", "SWOT-анализ")
        if not niche or not competitors:
            return {"statusCode": 400, "headers": headers, "body": json.dumps({"error": "Укажите нишу и конкурентов"})}
        system = "Ты — стратег по маркетингу и конкурентному анализу. Пиши конкретно, без воды."
        prompt = (
            f"Проведи анализ конкурентов: {analysis_type}\n"
            f"Ниша: {niche}\n"
            f"Конкуренты: {competitors}\n"
            f"Мои сильные стороны: {my_strengths}\n\n"
        )
        if analysis_type == "SWOT-анализ":
            prompt += "Дай полный SWOT-анализ относительно конкурентов. Для каждого конкурента — краткий разбор. Итог: мои возможности для отстройки."
        elif analysis_type == "Контент-стратегия":
            prompt += "Опиши предполагаемую контент-стратегию конкурентов и предложи 10 идей контента которые позволят обойти их по охватам и вовлечённости."
        elif analysis_type == "Ценовое позиционирование":
            prompt += "Проанализируй ценовое позиционирование в нише, предложи стратегию ценообразования и аргументацию ценности."
        elif analysis_type == "Аудитория":
            prompt += "Опиши аудиторию конкурентов, их боли и желания, предложи как привлечь эту аудиторию к себе."
        result = generate_text_with_openrouter(prompt, system)
        return {"statusCode": 200, "headers": headers, "body": json.dumps({"result": result, "type": "competitor-analysis"})}

    elif action == "case-generator":
        case_type = body.get("case_type", "Кейс (до/после)")
        client_name = body.get("client_name", "")
        problem = body.get("problem", "")
        solution = body.get("solution", "")
        result_text = body.get("result_text", "")
        metrics = body.get("metrics", "")
        style = body.get("style", "Деловой")
        if not problem or not solution:
            return {"statusCode": 400, "headers": headers, "body": json.dumps({"error": "Укажите проблему и решение"})}
        system = "Ты — профессиональный копирайтер. Создавай убедительные кейсы и отзывы которые продают. Русский язык."
        prompt = (
            f"Напиши '{case_type}' в стиле '{style}'.\n"
            f"Клиент: {client_name or 'клиент (обезличено)'}\n"
            f"Проблема/задача: {problem}\n"
            f"Решение: {solution}\n"
            f"Результат: {result_text}\n"
            f"Метрики: {metrics}\n\n"
            f"Требования: захватывающее начало, конкретные детали, эмоциональный отклик, "
            f"измеримые результаты, выводы которые убеждают потенциального клиента купить."
        )
        result = generate_text_with_openrouter(prompt, system)
        return {"statusCode": 200, "headers": headers, "body": json.dumps({"result": result, "type": "case-generator"})}

    elif action == "scenario":
        topic = body.get("topic", "")
        platform = body.get("platform", "Reels")
        duration = body.get("duration", "60 секунд")
        system = "Ты — сценарист вирусного контента для социальных сетей. Пиши на русском языке."
        prompt = f"Напиши детальный сценарий для {platform} на тему: '{topic}'. Длительность: {duration}. Укажи: крючок, основную часть, призыв к действию, советы по съёмке."
        result = generate_text_with_openrouter(prompt, system)
        return {"statusCode": 200, "headers": headers, "body": json.dumps({"result": result, "type": "scenario"})}

    elif action == "content-plan":
        niche = body.get("niche", "")
        period = body.get("period", "месяц")
        goals = body.get("goals", "рост подписчиков")
        platforms = body.get("platforms", ["Instagram"])
        platforms_str = ", ".join(platforms) if isinstance(platforms, list) else str(platforms)

        # Генерируем вступление
        intro_system = "Ты — стратег контент-маркетинга. Пиши структурированно, по делу, на русском языке."
        intro_prompt = (
            f"Напиши экспертное вступление к контент-плану для ниши '{niche}' на платформах: {platforms_str}. "
            f"Период: {period}. Цели: {goals}. "
            f"Структура (используй эти заголовки дословно):\n"
            f"1. Почему контент важен для этой ниши (3-4 предложения)\n"
            f"2. Особенности аудитории (3-4 предложения)\n"
            f"3. Что работает лучше всего (3-5 конкретных пунктов)\n"
            f"4. Ключевые метрики для отслеживания (4-5 метрик с пояснением)\n"
            f"5. Советы по публикациям (3-5 практических советов)\n"
            f"Пиши конкретно и практично, без воды."
        )
        intro_text = generate_text_with_openrouter(intro_prompt, intro_system)

        # Генерируем план
        plan_system = (
            "Ты — стратег контент-маркетинга. Отвечай ТОЛЬКО валидным JSON-массивом без лишнего текста, markdown и пояснений. "
            "Каждый элемент — объект со строго с полями: "
            "date (строка даты, например '01.06'), "
            "topic (развёрнутая тема поста — 1-2 предложения, конкретно и цепляюще), "
            "scenario (мини-сценарий: крючок, основная мысль, призыв — 2-3 предложения), "
            "format (один из: Reels/Карусель/Пост/Сторис/Статья/Видео/Shorts), "
            "notes (заметки по съёмке или оформлению), "
            "lifehacks (1-2 лайфхака для охвата и вовлечённости)."
        )
        plan_prompt = (
            f"Создай контент-план на {period} для ниши '{niche}'. "
            f"Платформы: {platforms_str}. Цели: {goals}. "
            f"Верни JSON-массив минимум из 20 объектов, равномерно по платформам. Только JSON."
        )
        raw_result = generate_text_with_openrouter(plan_prompt, plan_system)

        rows = None
        try:
            clean = raw_result.strip()
            if clean.startswith("```"):
                parts = clean.split("```")
                clean = parts[1] if len(parts) > 1 else clean
                if clean.startswith("json"):
                    clean = clean[4:]
                clean = clean.strip()
            rows = json.loads(clean)
        except Exception:
            rows = None

        return {"statusCode": 200, "headers": headers, "body": json.dumps({
            "result": raw_result,
            "intro": intro_text,
            "rows": rows,
            "type": "content_plan"
        })}

    elif action == "carousel":
        # Старый текстовый режим — оставляем для обратной совместимости
        topic = body.get("topic", "")
        slides_count = body.get("slides_count", 7)
        system = "Ты — создатель обучающего контента для Instagram. Пиши на русском языке."
        prompt = f"Создай структуру карусели из {slides_count} слайдов для Instagram на тему: '{topic}'. Для каждого слайда укажи: заголовок (до 8 слов), текст (до 30 слов), визуальное описание. Сделай первый слайд цепляющим."
        result = generate_text_with_openrouter(prompt, system)
        return {"statusCode": 200, "headers": headers, "body": json.dumps({"result": result, "type": "carousel"})}

    elif action == "carousel-structure":
        topic = body.get("topic", "")
        slides_count = body.get("slides_count", 7)
        ai_text = body.get("ai_text", True)
        if not topic:
            return {"statusCode": 400, "headers": headers, "body": json.dumps({"error": "Укажите тему карусели"})}
        try:
            slides = gemini_generate_slides(topic, slides_count, ai_text)
            return {"statusCode": 200, "headers": headers, "body": json.dumps({"slides": slides})}
        except Exception as e:
            print(f"[carousel-structure error] {e}")
            return {"statusCode": 500, "headers": headers, "body": json.dumps({"error": f"Ошибка генерации структуры: {str(e)}"})}

    elif action == "carousel-image":
        visual_prompt = body.get("visual_prompt", "")
        style_prompt = body.get("style_prompt", "minimalist clean design, soft colors")
        user_image_b64 = body.get("user_image_b64", None)
        slide_index = body.get("slide_index", 0)
        if not visual_prompt:
            return {"statusCode": 400, "headers": headers, "body": json.dumps({"error": "Укажите visual_prompt"})}
        try:
            if user_image_b64:
                img_bytes = generate_slide_image_with_photo(visual_prompt, style_prompt, user_image_b64)
            else:
                img_bytes = generate_slide_image_pollinations(visual_prompt, style_prompt)
            cdn_url = upload_image_to_s3(img_bytes, prefix="carousel")
            return {"statusCode": 200, "headers": headers, "body": json.dumps({
                "image_url": cdn_url,
                "slide_index": slide_index
            })}
        except Exception as e:
            print(f"[carousel-image error] slide {slide_index}: {e}")
            return {"statusCode": 500, "headers": headers, "body": json.dumps({"error": f"Ошибка генерации изображения: {str(e)}"})}


    elif action == "profile-analysis":
        niche = body.get("niche", "")
        followers = body.get("followers", 0)
        avg_likes = body.get("avg_likes", 0)
        system = "Ты — эксперт по продвижению в социальных сетях. Пиши на русском языке."
        prompt = f"Проанализируй Instagram профиль: ниша '{niche}', подписчиков: {followers}, среднее количество лайков: {avg_likes}. Дай подробный разбор профиля и пошаговую стратегию продвижения на 30 дней с конкретными действиями."
        result = generate_text_with_openrouter(prompt, system)
        return {"statusCode": 200, "headers": headers, "body": json.dumps({"result": result, "type": "profile_analysis"})}

    elif action == "funnel":
        platform = body.get("platform", "Instagram")
        product = body.get("product", "")
        audience = body.get("audience", "")
        system = "Ты — эксперт по продажам и воронкам в социальных сетях. Пиши на русском языке."
        prompt = f"Создай продающую воронку для {platform}: продукт/услуга: '{product}', целевая аудитория: '{audience}'. Укажи: все этапы воронки, типы контента для каждого этапа, скрипты сообщений, точки касания."
        result = generate_text_with_openrouter(prompt, system)
        return {"statusCode": 200, "headers": headers, "body": json.dumps({"result": result, "type": "funnel"})}

    elif action == "presentation":
        topic = body.get("topic", "")
        slides_count = body.get("slides_count", 10)
        theme = body.get("theme", "dark")
        custom_style = body.get("custom_style", "")  # пользовательский стиль картинок

        # 1. Генерируем структуру через OpenRouter
        system = "Ты — профессиональный презентационный дизайнер и эксперт в теме. Отвечай ТОЛЬКО валидным JSON-массивом без markdown и пояснений."
        prompt = (
            f"Создай структуру презентации из {slides_count} слайдов на тему: '{topic}'.\n"
            f"Для каждого слайда верни объект строго:\n"
            f'{{"slide": 1, "title": "Заголовок слайда", "bullets": ["Развёрнутый тезис с конкретными данными или советом", "Второй тезис — практичный и понятный", "Третий тезис — итог или следующий шаг"], "image_query": "very specific english prompt for image generation matching slide topic, 8-12 words, photorealistic"}}\n\n'
            f"Требования:\n"
            f"- Первый слайд — красивая обложка, bullets[0] = краткий подзаголовок презентации\n"
            f"- Последний слайд — выводы и призыв к действию\n"
            f"- Каждый тезис — 1-2 предложения, конкретно, без воды, с пользой\n"
            f"- image_query: ОЧЕНЬ конкретный запрос строго по теме этого слайда (не общий фон!). Например для слайда о SEO: 'SEO analytics dashboard computer screen data charts'\n"
            f"- Только JSON-массив, без markdown."
        )
        raw = generate_text_with_openrouter(prompt, system)
        slides_data = None
        try:
            clean = raw.strip()
            if clean.startswith("```"):
                parts = clean.split("```")
                clean = parts[1] if len(parts) > 1 else clean
                if clean.startswith("json"): clean = clean[4:]
                clean = clean.strip()
            slides_data = json.loads(clean)
        except Exception:
            slides_data = None

        if not slides_data:
            return {"statusCode": 200, "headers": headers, "body": json.dumps({"result": raw, "type": "presentation"})}

        # 2. Темы оформления
        THEMES = {
            "dark":      {"bg": (15, 15, 25),    "accent": (99, 102, 241),  "title": (255,255,255), "body": (180,180,200), "sub": (120,120,150)},
            "light":     {"bg": (248, 249, 252),  "accent": (99, 102, 241),  "title": (20, 20, 40),  "body": (60, 60, 80),  "sub": (120,120,140)},
            "corporate": {"bg": (10, 20, 50),     "accent": (0, 180, 216),   "title": (255,255,255), "body": (180,210,240), "sub": (100,140,180)},
            "creative":  {"bg": (20, 10, 35),     "accent": (236, 72, 153),  "title": (255,255,255), "body": (220,180,240), "sub": (160,120,200)},
            "minimal":   {"bg": (255, 255, 255),  "accent": (30, 30, 30),    "title": (10, 10, 10),  "body": (60, 60, 60),  "sub": (140,140,140)},
        }
        T = THEMES.get(theme, THEMES["dark"])

        # 3. Строим PPTX
        from pptx import Presentation as PPTXPresentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        import io as _io

        prs = PPTXPresentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)

        def rgb(r,g,b): return RGBColor(r,g,b)
        def add_rect(slide, l, t, w, h, fill_rgb, alpha=None):
            shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
            shape.fill.solid()
            shape.fill.fore_color.rgb = rgb(*fill_rgb)
            shape.line.fill.background()
            return shape

        def add_text(slide, text, l, t, w, h, font_size, color_rgb, bold=False, align=PP_ALIGN.LEFT):
            txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = align
            run = p.add_run()
            run.text = text
            run.font.size = Pt(font_size)
            run.font.color.rgb = rgb(*color_rgb)
            run.font.bold = bold
            return txBox

        def fetch_image(query: str) -> bytes | None:
            try:
                import urllib.parse as _up
                # Если пользователь задал свой стиль — добавляем его
                style_suffix = f", {custom_style}" if custom_style else ", professional photography, high quality, sharp focus"
                full_query = query + style_suffix
                enc = _up.quote(full_query)
                url = f"https://image.pollinations.ai/prompt/{enc}?width=1920&height=1080&model=flux&nologo=true&enhance=false"
                req2 = urllib.request.Request(url, method="GET")
                req2.add_header("User-Agent", "Mozilla/5.0")
                with urllib.request.urlopen(req2, timeout=45) as r:
                    return r.read()
            except Exception as ex:
                print(f"[pptx img error] {ex}")
                return None

        for idx, slide_info in enumerate(slides_data):
            slide_layout = prs.slide_layouts[6]  # blank
            slide = prs.slides.add_slide(slide_layout)
            title_text   = slide_info.get("title", f"Слайд {idx+1}")
            bullets      = slide_info.get("bullets", [])
            image_query  = slide_info.get("image_query", topic)
            is_cover     = (idx == 0)

            # Фон
            add_rect(slide, 0, 0, 13.33, 7.5, T["bg"])

            # Картинка
            img_bytes = fetch_image(image_query)
            if img_bytes:
                try:
                    from PIL import Image as PILImage
                    pil = PILImage.open(_io.BytesIO(img_bytes)).convert("RGB")
                    if is_cover:
                        # Обложка — полная картинка с тёмным оверлеем
                        pil = pil.resize((1920, 1080), PILImage.LANCZOS)
                        buf = _io.BytesIO(); pil.save(buf, "PNG"); buf.seek(0)
                        slide.shapes.add_picture(buf, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
                        # Тёмный оверлей
                        overlay = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
                        overlay.fill.solid()
                        overlay.fill.fore_color.rgb = rgb(0,0,0)
                        overlay.fill.fore_color.theme_color = None
                        from pptx.util import Pt as _Pt
                        overlay.line.fill.background()
                        # прозрачность через XML
                        from lxml import etree
                        sp_el = overlay.element
                        sp_pr = sp_el.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
                        if sp_pr is not None:
                            srgb = sp_pr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
                            if srgb is None:
                                srgb = etree.SubElement(sp_pr, '{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
                                srgb.set('val', '000000')
                            alpha_el = etree.SubElement(srgb, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
                            alpha_el.set('val', '60000')
                    else:
                        # Обычный слайд — картинка справа
                        pil = pil.resize((760, 540), PILImage.LANCZOS)
                        buf = _io.BytesIO(); pil.save(buf, "PNG"); buf.seek(0)
                        slide.shapes.add_picture(buf, Inches(7.8), Inches(0.8), Inches(5.0), Inches(5.9))
                        # Акцентная полоска слева
                        add_rect(slide, 0, 0, 0.08, 7.5, T["accent"])
                except Exception as ex:
                    print(f"[pptx img render error] {ex}")

            if is_cover:
                # Номер слайда / надпись сверху
                add_text(slide, topic.upper(), 1.0, 1.2, 11.0, 0.5, 11, T["sub"], align=PP_ALIGN.CENTER)
                # Большой заголовок
                add_text(slide, title_text, 0.8, 2.2, 11.7, 2.0, 44, T["title"], bold=True, align=PP_ALIGN.CENTER)
                # Подзаголовок
                if bullets:
                    add_text(slide, bullets[0] if isinstance(bullets, list) else str(bullets), 1.0, 4.4, 11.3, 0.8, 20, T["body"], align=PP_ALIGN.CENTER)
                # Акцентная линия
                add_rect(slide, 5.0, 4.1, 3.33, 0.06, T["accent"])
            else:
                # Заголовок
                add_text(slide, title_text, 0.4, 0.3, 7.0, 1.0, 28, T["title"], bold=True)
                # Акцентная линия под заголовком
                add_rect(slide, 0.4, 1.35, 2.5, 0.05, T["accent"])
                # Буллеты
                y = 1.6
                for bullet in (bullets if isinstance(bullets, list) else []):
                    bullet_str = f"• {bullet}"
                    add_text(slide, bullet_str, 0.55, y, 6.9, 0.75, 16, T["body"])
                    y += 0.85
                # Номер слайда
                add_text(slide, str(idx + 1), 12.5, 6.9, 0.5, 0.4, 11, T["sub"], align=PP_ALIGN.RIGHT)

        # 4. Сохраняем и загружаем в S3
        buf = _io.BytesIO()
        prs.save(buf)
        pptx_bytes = buf.getvalue()
        cdn_url = upload_image_to_s3(pptx_bytes, prefix="presentations")
        # Меняем расширение в URL
        cdn_url = cdn_url.replace(".png", ".pptx")
        # Загружаем с правильным типом
        s3 = get_s3()
        key = f"images/presentations/{uuid.uuid4()}.pptx"
        s3.put_object(Bucket="files", Key=key, Body=pptx_bytes, ContentType="application/vnd.openxmlformats-officedocument.presentationml.presentation")
        access_key = os.environ["AWS_ACCESS_KEY_ID"]
        cdn_url = f"https://cdn.poehali.dev/projects/{access_key}/bucket/{key}"

        return {"statusCode": 200, "headers": headers, "body": json.dumps({
            "result": raw,
            "pptx_url": cdn_url,
            "slides_count": len(slides_data),
            "type": "presentation"
        })}

    elif action == "guide":
        topic = body.get("topic", "")
        guide_type = body.get("guide_type", "пошаговый")
        system = "Ты — эксперт-методолог. Пиши на русском языке."
        prompt = f"Создай {guide_type} гайд на тему: '{topic}'. Структура: введение, пошаговые инструкции с деталями, советы, частые ошибки, заключение."
        result = generate_text_with_openrouter(prompt, system)
        return {"statusCode": 200, "headers": headers, "body": json.dumps({"result": result, "type": "guide"})}

    elif action == "product-card":
        product_name = body.get("product_name", "")
        features = body.get("features", "")
        price = body.get("price", "")
        system = "Ты — маркетолог маркетплейсов. Пиши на русском языке."
        prompt = f"Создай продающую карточку товара для маркетплейса: '{product_name}', характеристики: {features}, цена: {price}. Напиши: заголовок (с ключевыми словами), описание, bullet-points с преимуществами, SEO-теги."
        result = generate_text_with_openrouter(prompt, system)
        return {"statusCode": 200, "headers": headers, "body": json.dumps({"result": result, "type": "product_card"})}

    elif action == "reels-analysis":
        video_description = body.get("video_description", "")
        views = body.get("views", 0)
        likes = body.get("likes", 0)
        comments = body.get("comments", 0)
        system = "Ты — аналитик видео-контента. Пиши на русском языке."
        prompt = f"Проанализируй видео Reels: описание: '{video_description}', просмотры: {views}, лайки: {likes}, комментарии: {comments}. Дай анализ: вовлечённость, что работает, что улучшить, рекомендации по продвижению."
        result = generate_text_with_openrouter(prompt, system)
        return {"statusCode": 200, "headers": headers, "body": json.dumps({"result": result, "type": "reels_analysis"})}

    return {"statusCode": 404, "headers": headers, "body": json.dumps({"error": "Действие не найдено"})}